"""PowerPoint parser using python-pptx for native parsing."""

from pathlib import Path

import structlog
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.table import Table

from rag_assist.ingestion.models import (
    DocumentMetadata,
    DocumentType,
    ParsedDocument,
    SlideContent,
    TextBlock,
)
from rag_assist.ingestion.parsers.base_parser import BaseParser

logger = structlog.get_logger(__name__)


class PPTXParser(BaseParser):
    """Native PowerPoint parser preserving slide structure.

    Extracts:
    - Slide titles
    - Body text (bullets, paragraphs)
    - Speaker notes
    - Tables
    - SmartArt text
    """

    supported_types = [DocumentType.PPTX]

    def __init__(self, include_speaker_notes: bool = True):
        """Initialize PPTX parser.

        Args:
            include_speaker_notes: Whether to include speaker notes in extraction.
        """
        self.include_speaker_notes = include_speaker_notes

    def can_parse(self, file_path: str | Path) -> bool:
        """Check if file is a PowerPoint."""
        path = Path(file_path)
        return path.suffix.lower() in [".pptx", ".ppt"]

    def parse(self, file_path: str | Path, metadata: DocumentMetadata) -> ParsedDocument:
        """Parse PowerPoint document.

        Args:
            file_path: Path to PPTX file.
            metadata: Document metadata.

        Returns:
            ParsedDocument with extracted slides.
        """
        self.log_parsing_start(file_path)
        path = Path(file_path)
        slides: list[SlideContent] = []
        errors: list[str] = []

        try:
            prs = Presentation(str(path))
            metadata.total_slides = len(prs.slides)

            for slide_num, slide in enumerate(prs.slides, 1):
                try:
                    slide_content = self._extract_slide(slide, slide_num)
                    slides.append(slide_content)
                except Exception as e:
                    error_msg = f"Error parsing slide {slide_num}: {str(e)}"
                    errors.append(error_msg)
                    logger.warning(error_msg)
                    # Add empty slide to maintain numbering
                    slides.append(
                        SlideContent(
                            slide_number=slide_num,
                            title=f"[Error parsing slide {slide_num}]",
                        )
                    )

            metadata.extraction_method = "native"

        except Exception as e:
            self.log_parsing_error(file_path, e)
            errors.append(f"Failed to open PowerPoint: {str(e)}")

        result = ParsedDocument(
            metadata=metadata,
            content=slides,
            extraction_method="native",
            processing_errors=errors,
        )

        self.log_parsing_complete(file_path, result)
        return result

    def _extract_slide(self, slide, slide_number: int) -> SlideContent:
        """Extract content from a single slide.

        Args:
            slide: python-pptx slide object.
            slide_number: Slide number (1-indexed).

        Returns:
            SlideContent with extracted text.
        """
        # Extract title
        title = self._extract_title(slide)

        # Extract body text from all shapes
        body_text = self._extract_body_text(slide)

        # Extract speaker notes
        speaker_notes = ""
        if self.include_speaker_notes:
            speaker_notes = self._extract_speaker_notes(slide)

        # Extract tables
        tables = self._extract_tables(slide)

        # Get layout type
        layout_type = ""
        try:
            if slide.slide_layout:
                layout_type = slide.slide_layout.name
        except Exception:
            pass

        return SlideContent(
            slide_number=slide_number,
            title=title,
            body_text=body_text,
            speaker_notes=speaker_notes,
            tables=tables,
            layout_type=layout_type,
        )

    def _extract_title(self, slide) -> str:
        """Extract slide title.

        Args:
            slide: python-pptx slide object.

        Returns:
            Slide title text.
        """
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            return self.normalize_text(slide.shapes.title.text)
        return ""

    def _extract_body_text(self, slide) -> list[TextBlock]:
        """Extract body text from all shapes.

        Args:
            slide: python-pptx slide object.

        Returns:
            List of TextBlocks with extracted text.
        """
        text_blocks = []

        for shape in slide.shapes:
            # Skip title shape (already extracted)
            if shape == slide.shapes.title:
                continue

            # Handle grouped shapes
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                text_blocks.extend(self._extract_from_group(shape))
                continue

            # Handle text frames
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = self.normalize_text(paragraph.text)
                    if text:
                        text_blocks.append(
                            TextBlock(
                                text=text,
                                level=paragraph.level,
                                font_size=self._get_font_size(paragraph),
                                is_bold=self._is_bold(paragraph),
                                is_heading=paragraph.level == 0,
                                block_type="list_item" if paragraph.level > 0 else "paragraph",
                            )
                        )

        return text_blocks

    def _extract_from_group(self, group_shape) -> list[TextBlock]:
        """Extract text from grouped shapes.

        Args:
            group_shape: python-pptx group shape object.

        Returns:
            List of TextBlocks from grouped shapes.
        """
        text_blocks = []

        try:
            for shape in group_shape.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = self.normalize_text(paragraph.text)
                        if text:
                            text_blocks.append(
                                TextBlock(
                                    text=text,
                                    level=paragraph.level,
                                    block_type="paragraph",
                                )
                            )
        except Exception:
            pass

        return text_blocks

    def _extract_speaker_notes(self, slide) -> str:
        """Extract speaker notes from slide.

        Args:
            slide: python-pptx slide object.

        Returns:
            Speaker notes text.
        """
        try:
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_frame = notes_slide.notes_text_frame
                if notes_frame:
                    return self.normalize_text(notes_frame.text)
        except Exception:
            pass
        return ""

    def _extract_tables(self, slide) -> list[list[list[str]]]:
        """Extract tables from slide.

        Args:
            slide: python-pptx slide object.

        Returns:
            List of tables, each table is a list of rows, each row is a list of cells.
        """
        tables = []

        for shape in slide.shapes:
            if shape.has_table:
                table: Table = shape.table
                table_data = []

                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        cell_text = self.normalize_text(cell.text)
                        row_data.append(cell_text)
                    table_data.append(row_data)

                if table_data:
                    tables.append(table_data)

        return tables

    def _get_font_size(self, paragraph) -> float | None:
        """Get font size from paragraph.

        Args:
            paragraph: python-pptx paragraph object.

        Returns:
            Font size in points or None.
        """
        try:
            for run in paragraph.runs:
                if run.font.size:
                    return run.font.size.pt
        except Exception:
            pass
        return None

    def _is_bold(self, paragraph) -> bool:
        """Check if paragraph is bold.

        Args:
            paragraph: python-pptx paragraph object.

        Returns:
            True if paragraph has bold formatting.
        """
        try:
            for run in paragraph.runs:
                if run.font.bold:
                    return True
        except Exception:
            pass
        return False
