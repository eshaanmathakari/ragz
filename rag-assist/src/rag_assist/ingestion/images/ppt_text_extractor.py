"""Step 1: PPT Text Extraction using python-pptx.

This module extracts text content from PowerPoint presentations:
- Slide titles
- Body text (bullets, paragraphs)
- Speaker notes
- Tables
- SmartArt text (from grouped shapes)

Also detects which slides contain visual content (images, SmartArt, charts).
"""

from __future__ import annotations

import hashlib
import logging
import uuid
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from models import PPTExtractionResult, SlideTextContent

# Try structlog, fallback to standard logging
try:
    import structlog
    logger = structlog.get_logger(__name__)
except ImportError:
    logging.basicConfig(level=logging.INFO)
    logger = logging.getLogger(__name__)


class PPTXTextExtractor:
    """Extract text content from PPTX slides.

    Usage:
        extractor = PPTXTextExtractor()
        result = extractor.extract("presentation.pptx")

        for slide in result.slides:
            print(f"Slide {slide.slide_number}: {slide.title}")
            print(slide.full_text)
    """

    def __init__(
        self,
        include_speaker_notes: bool = True,
        include_tables: bool = True,
        detect_visual_content: bool = True,
        smartart_shape_threshold: int = 5,
    ):
        """Initialize extractor.

        Args:
            include_speaker_notes: Extract speaker notes.
            include_tables: Extract table content as text.
            detect_visual_content: Flag slides with images/SmartArt/charts.
            smartart_shape_threshold: Min shapes in group to consider as SmartArt.
        """
        self.include_speaker_notes = include_speaker_notes
        self.include_tables = include_tables
        self.detect_visual_content = detect_visual_content
        self.smartart_shape_threshold = smartart_shape_threshold

    def extract(self, pptx_path: str | Path) -> PPTExtractionResult:
        """Extract text from all slides.

        Args:
            pptx_path: Path to PPTX file.

        Returns:
            PPTExtractionResult with all slide content.
        """
        path = Path(pptx_path)
        filename = path.name
        document_id = self._generate_document_id(path)
        slides: list[SlideTextContent] = []
        errors: list[str] = []

        logger.info("Extracting text from PPTX", filename=filename)

        try:
            prs = Presentation(str(path))
            total_slides = len(prs.slides)

            for slide_num, slide in enumerate(prs.slides, 1):
                try:
                    slide_content = self._extract_slide(
                        slide, slide_num, document_id, filename
                    )
                    slides.append(slide_content)
                except Exception as e:
                    error_msg = f"Error parsing slide {slide_num}: {str(e)}"
                    errors.append(error_msg)
                    logger.warning(error_msg)
                    # Add minimal slide to maintain numbering
                    slides.append(
                        SlideTextContent(
                            document_id=document_id,
                            filename=filename,
                            slide_number=slide_num,
                            title=f"[Error parsing slide {slide_num}]",
                            body_text="",
                        )
                    )

            logger.info(
                "PPTX text extraction complete",
                filename=filename,
                total_slides=total_slides,
                slides_extracted=len(slides),
                errors=len(errors),
            )

        except Exception as e:
            error_msg = f"Failed to open PowerPoint: {str(e)}"
            errors.append(error_msg)
            logger.error(error_msg)
            total_slides = 0

        return PPTExtractionResult(
            document_id=document_id,
            filename=filename,
            total_slides=total_slides,
            slides=slides,
            errors=errors,
        )

    def extract_slide_numbers_with_visuals(self, pptx_path: str | Path) -> list[int]:
        """Get slide numbers that contain visual content.

        Args:
            pptx_path: Path to PPTX file.

        Returns:
            List of slide numbers (1-indexed) with images/SmartArt/charts.
        """
        result = self.extract(pptx_path)
        return result.slides_with_visuals

    def _extract_slide(
        self, slide, slide_num: int, document_id: str, filename: str
    ) -> SlideTextContent:
        """Extract content from a single slide."""
        # Extract title
        title = self._extract_title(slide)

        # Extract body text
        body_text = self._extract_body_text(slide)

        # Extract speaker notes
        speaker_notes = ""
        if self.include_speaker_notes:
            speaker_notes = self._extract_speaker_notes(slide)

        # Extract tables
        tables: list[str] = []
        if self.include_tables:
            tables = self._extract_tables(slide)

        # Detect visual content
        has_visual = False
        if self.detect_visual_content:
            has_visual = self._has_visual_content(slide)

        return SlideTextContent(
            document_id=document_id,
            filename=filename,
            slide_number=slide_num,
            title=title,
            body_text=body_text,
            speaker_notes=speaker_notes,
            tables=tables,
            has_visual_content=has_visual,
        )

    def _extract_title(self, slide) -> str:
        """Extract slide title."""
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            return self._normalize_text(slide.shapes.title.text)
        return ""

    def _extract_body_text(self, slide) -> str:
        """Extract body text from all shapes."""
        text_parts = []

        for shape in slide.shapes:
            # Skip title shape (already extracted)
            if shape == slide.shapes.title:
                continue

            # Handle grouped shapes (including SmartArt)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_text = self._extract_from_group(shape)
                if group_text:
                    text_parts.append(group_text)
                continue

            # Handle text frames
            if shape.has_text_frame:
                frame_text = self._extract_text_frame(shape.text_frame)
                if frame_text:
                    text_parts.append(frame_text)

        return "\n".join(text_parts)

    def _extract_from_group(self, group_shape) -> str:
        """Extract text from grouped shapes (SmartArt, etc.)."""
        text_parts = []

        try:
            for shape in group_shape.shapes:
                if shape.has_text_frame:
                    frame_text = self._extract_text_frame(shape.text_frame)
                    if frame_text:
                        text_parts.append(frame_text)
        except Exception:
            pass

        return "\n".join(text_parts)

    def _extract_text_frame(self, text_frame) -> str:
        """Extract text from a text frame."""
        paragraphs = []
        for paragraph in text_frame.paragraphs:
            text = self._normalize_text(paragraph.text)
            if text:
                # Add bullet prefix for indented items
                if paragraph.level > 0:
                    text = "  " * paragraph.level + "• " + text
                paragraphs.append(text)
        return "\n".join(paragraphs)

    def _extract_speaker_notes(self, slide) -> str:
        """Extract speaker notes from slide."""
        try:
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_frame = notes_slide.notes_text_frame
                if notes_frame:
                    return self._normalize_text(notes_frame.text)
        except Exception:
            pass
        return ""

    def _extract_tables(self, slide) -> list[str]:
        """Extract tables from slide as text."""
        tables = []

        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                rows_text = []

                for row in table.rows:
                    cells = []
                    for cell in row.cells:
                        cell_text = self._normalize_text(cell.text)
                        cells.append(cell_text)
                    rows_text.append(" | ".join(cells))

                if rows_text:
                    tables.append("\n".join(rows_text))

        return tables

    def _has_visual_content(self, slide) -> bool:
        """Check if slide has visual content (images, SmartArt, charts)."""
        for shape in slide.shapes:
            # Check for pictures
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return True

            # Check for charts
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                return True

            # Check for SmartArt (grouped shapes with many children)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                try:
                    if hasattr(shape, "shapes"):
                        if len(shape.shapes) >= self.smartart_shape_threshold:
                            return True
                except Exception:
                    pass

            # Check for diagram namespace in XML
            try:
                element = shape._element
                xml_str = element.xml if hasattr(element, "xml") else str(element)
                if "dgm:" in xml_str or "diagram" in xml_str.lower():
                    return True
            except Exception:
                pass

        return False

    def _normalize_text(self, text: str) -> str:
        """Normalize whitespace in text."""
        if not text:
            return ""
        # Replace multiple whitespace with single space
        import re
        text = re.sub(r"\s+", " ", text)
        return text.strip()

    def _generate_document_id(self, path: Path) -> str:
        """Generate document ID from file path and content hash."""
        try:
            content = path.read_bytes()
            hash_digest = hashlib.sha256(content).hexdigest()[:12]
            return f"{path.stem}_{hash_digest}"
        except Exception:
            return f"{path.stem}_{uuid.uuid4().hex[:12]}"
