"""PPTX Image Extraction Module (v2 - PDF Approach).

A portable module for extracting images from PowerPoint presentations:
1. Detects slides with SmartArt/images using python-pptx
2. Converts PPTX to PDF using LibreOffice (renders SmartArt properly)
3. Extracts specific PDF pages as PNG images using PyMuPDF
4. Uploads to S3 with slide number mapping
"""

from __future__ import annotations

from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
import hashlib
import uuid
import subprocess
import tempfile
import logging

# Try structlog, fallback to standard logging
try:
    import structlog
    logger = structlog.get_logger(__name__)
except ImportError:
    logging.basicConfig(level=logging.INFO)
    logger = logging.getLogger(__name__)

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# =============================================================================
# Data Classes
# =============================================================================


@dataclass
class ImageInfo:
    """Information about an extracted image."""

    image_id: str
    document_id: str
    filename: str
    slide_number: int
    image_type: str  # "pdf_render", "picture", "smartart", "chart"
    s3_uri: str = ""
    s3_key: str = ""
    width_px: int = 0
    height_px: int = 0
    size_bytes: int = 0
    extracted_at: datetime = field(default_factory=datetime.utcnow)


@dataclass
class SlideImageMapping:
    """Mapping of visual content found on a slide."""

    document_id: str
    filename: str
    slide_number: int
    slide_title: str = ""
    has_pictures: bool = False
    has_smartart: bool = False
    has_charts: bool = False
    images: list[ImageInfo] = field(default_factory=list)


@dataclass
class ExtractionResult:
    """Result of the full extraction pipeline."""

    document_id: str
    filename: str
    total_slides: int
    slides_with_images: int
    images: list[ImageInfo] = field(default_factory=list)
    mappings: list[SlideImageMapping] = field(default_factory=list)
    errors: list[str] = field(default_factory=list)


# =============================================================================
# PPTXImageDetector - Detect slides with visual content
# =============================================================================


class PPTXImageDetector:
    """Detect slides containing images, SmartArt, and charts."""

    def __init__(
        self,
        include_pictures: bool = True,
        include_smartart: bool = True,
        include_charts: bool = True,
        smartart_shape_threshold: int = 5,
    ):
        """Initialize detector.

        Args:
            include_pictures: Detect embedded pictures.
            include_smartart: Detect SmartArt diagrams.
            include_charts: Detect charts.
            smartart_shape_threshold: Min shapes in group to consider as SmartArt.
        """
        self.include_pictures = include_pictures
        self.include_smartart = include_smartart
        self.include_charts = include_charts
        self.smartart_shape_threshold = smartart_shape_threshold

    def detect_image_slides(
        self, pptx_path: str | Path
    ) -> list[SlideImageMapping]:
        """Scan PPTX and return slides that contain visual content.

        Args:
            pptx_path: Path to PPTX file.

        Returns:
            List of SlideImageMapping for slides with images.
        """
        path = Path(pptx_path)
        document_id = self._generate_document_id(path)
        filename = path.name
        mappings = []

        try:
            prs = Presentation(str(path))

            for slide_num, slide in enumerate(prs.slides, 1):
                mapping = self._analyze_slide(slide, slide_num, document_id, filename)
                if mapping:
                    mappings.append(mapping)

            logger.info(
                "Detected slides with images",
                filename=filename,
                total_slides=len(prs.slides),
                slides_with_images=len(mappings),
            )

        except Exception as e:
            logger.error("Failed to detect images", filename=filename, error=str(e))
            raise

        return mappings

    def get_slide_numbers_with_images(self, pptx_path: str | Path) -> list[int]:
        """Get just the slide numbers that have visual content.

        Args:
            pptx_path: Path to PPTX file.

        Returns:
            List of slide numbers (1-indexed).
        """
        mappings = self.detect_image_slides(pptx_path)
        return [m.slide_number for m in mappings]

    def _analyze_slide(
        self, slide, slide_num: int, document_id: str, filename: str
    ) -> SlideImageMapping | None:
        """Analyze a single slide for visual content."""
        has_pictures = False
        has_smartart = False
        has_charts = False

        for shape in slide.shapes:
            if self.include_pictures and self._is_picture(shape):
                has_pictures = True
            if self.include_smartart and self._is_smartart(shape):
                has_smartart = True
            if self.include_charts and self._is_chart(shape):
                has_charts = True

        if not (has_pictures or has_smartart or has_charts):
            return None

        # Get slide title
        title = ""
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text.strip()

        return SlideImageMapping(
            document_id=document_id,
            filename=filename,
            slide_number=slide_num,
            slide_title=title,
            has_pictures=has_pictures,
            has_smartart=has_smartart,
            has_charts=has_charts,
        )

    def _is_picture(self, shape) -> bool:
        """Check if shape is an embedded picture."""
        return shape.shape_type == MSO_SHAPE_TYPE.PICTURE

    def _is_smartart(self, shape) -> bool:
        """Check if shape is SmartArt."""
        # Check for grouped shapes with many children (SmartArt heuristic)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                if hasattr(shape, "shapes"):
                    return len(shape.shapes) >= self.smartart_shape_threshold
            except Exception:
                pass

        # Check for GraphicFrame with diagram namespace
        try:
            element = shape._element
            xml_str = element.xml if hasattr(element, "xml") else str(element)
            if "dgm:" in xml_str or "diagram" in xml_str.lower():
                return True
        except Exception:
            pass

        return False

    def _is_chart(self, shape) -> bool:
        """Check if shape is a chart."""
        return shape.shape_type == MSO_SHAPE_TYPE.CHART

    def _generate_document_id(self, path: Path) -> str:
        """Generate document ID from file path and content hash."""
        try:
            content = path.read_bytes()
            hash_digest = hashlib.sha256(content).hexdigest()[:12]
            return f"{path.stem}_{hash_digest}"
        except Exception:
            return f"{path.stem}_{uuid.uuid4().hex[:12]}"


# =============================================================================
# PPTXtoPDFConverter - Convert PPTX to PDF using LibreOffice
# =============================================================================


class PPTXtoPDFConverter:
    """Convert PPTX to PDF using LibreOffice headless."""

    def __init__(self, timeout: int = 120):
        """Initialize converter.

        Args:
            timeout: Conversion timeout in seconds.
        """
        self.timeout = timeout
        self._libreoffice_available = None
        self._libreoffice_cmd = None

    def convert(self, pptx_path: str | Path, output_dir: str | Path | None = None) -> Path:
        """Convert PPTX to PDF.

        Args:
            pptx_path: Path to PPTX file.
            output_dir: Output directory (defaults to temp dir).

        Returns:
            Path to generated PDF file.

        Raises:
            RuntimeError: If LibreOffice not available or conversion fails.
        """
        pptx_path = Path(pptx_path)

        if not self._check_libreoffice():
            raise RuntimeError(
                "LibreOffice not available. Install with: "
                "brew install libreoffice (macOS) or "
                "sudo apt install libreoffice (Linux) or "
                "sudo yum install libreoffice-headless (Amazon Linux)"
            )

        # Use temp dir if no output dir specified
        if output_dir is None:
            output_dir = Path(tempfile.mkdtemp(prefix="pptx_to_pdf_"))
        else:
            output_dir = Path(output_dir)
            output_dir.mkdir(parents=True, exist_ok=True)

        # Run LibreOffice conversion
        cmd = [
            self._libreoffice_cmd,
            "--headless",
            "--convert-to", "pdf",
            "--outdir", str(output_dir),
            str(pptx_path),
        ]

        logger.info("Converting PPTX to PDF", pptx=str(pptx_path), output_dir=str(output_dir))

        try:
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=self.timeout,
            )

            if result.returncode != 0:
                raise RuntimeError(f"LibreOffice conversion failed: {result.stderr}")

        except subprocess.TimeoutExpired:
            raise RuntimeError(f"Conversion timed out after {self.timeout}s")

        # Find the generated PDF
        pdf_name = pptx_path.stem + ".pdf"
        pdf_path = output_dir / pdf_name

        if not pdf_path.exists():
            # Try to find any PDF in output dir
            pdf_files = list(output_dir.glob("*.pdf"))
            if pdf_files:
                pdf_path = pdf_files[0]
            else:
                raise RuntimeError(f"PDF not generated. Expected: {pdf_path}")

        logger.info("PDF conversion complete", pdf=str(pdf_path))
        return pdf_path

    def _check_libreoffice(self) -> bool:
        """Check if LibreOffice is available."""
        if self._libreoffice_available is not None:
            return self._libreoffice_available

        # Try different command names
        for cmd in ["soffice", "libreoffice", "/Applications/LibreOffice.app/Contents/MacOS/soffice"]:
            try:
                result = subprocess.run(
                    [cmd, "--version"],
                    capture_output=True,
                    text=True,
                    timeout=5,
                )
                if result.returncode == 0:
                    self._libreoffice_available = True
                    self._libreoffice_cmd = cmd
                    logger.info("LibreOffice found", cmd=cmd)
                    return True
            except (subprocess.SubprocessError, FileNotFoundError):
                continue

        self._libreoffice_available = False
        logger.warning("LibreOffice not found")
        return False


# =============================================================================
# PDFPageExtractor - Extract specific PDF pages as PNG images
# =============================================================================


class PDFPageExtractor:
    """Extract specific PDF pages as PNG images using PyMuPDF."""

    def __init__(self, dpi: int = 150):
        """Initialize extractor.

        Args:
            dpi: Resolution for rendering (higher = better quality, larger files).
        """
        self.dpi = dpi
        self.zoom = dpi / 72  # PDF default is 72 DPI

    def extract_pages(
        self,
        pdf_path: str | Path,
        page_numbers: list[int],
        document_id: str,
        filename: str = "",
    ) -> list[tuple[int, bytes, ImageInfo]]:
        """Render specific PDF pages as PNG images.

        Args:
            pdf_path: Path to PDF file.
            page_numbers: Page numbers to extract (1-indexed).
            document_id: Document identifier for metadata.
            filename: Original filename for metadata.

        Returns:
            List of (page_number, png_bytes, ImageInfo).
        """
        import fitz  # PyMuPDF

        pdf_path = Path(pdf_path)
        if not filename:
            filename = pdf_path.stem

        results = []

        try:
            doc = fitz.open(str(pdf_path))
            total_pages = len(doc)

            for page_num in page_numbers:
                if page_num < 1 or page_num > total_pages:
                    logger.warning(
                        "Invalid page number",
                        page=page_num,
                        total=total_pages,
                    )
                    continue

                try:
                    # Get page (0-indexed)
                    page = doc[page_num - 1]

                    # Render to pixmap
                    mat = fitz.Matrix(self.zoom, self.zoom)
                    pix = page.get_pixmap(matrix=mat)

                    # Convert to PNG bytes
                    img_bytes = pix.tobytes("png")

                    # Create ImageInfo
                    info = ImageInfo(
                        image_id=f"{document_id}_slide_{page_num:03d}",
                        document_id=document_id,
                        filename=filename,
                        slide_number=page_num,
                        image_type="pdf_render",
                        width_px=pix.width,
                        height_px=pix.height,
                        size_bytes=len(img_bytes),
                    )

                    results.append((page_num, img_bytes, info))

                    logger.info(
                        "Extracted page",
                        page=page_num,
                        width=pix.width,
                        height=pix.height,
                        size_kb=len(img_bytes) // 1024,
                    )

                except Exception as e:
                    logger.error("Failed to extract page", page=page_num, error=str(e))

            doc.close()

        except Exception as e:
            logger.error("Failed to open PDF", path=str(pdf_path), error=str(e))
            raise

        return results

    def extract_all_pages(
        self,
        pdf_path: str | Path,
        document_id: str,
        filename: str = "",
    ) -> list[tuple[int, bytes, ImageInfo]]:
        """Extract all pages from PDF.

        Args:
            pdf_path: Path to PDF file.
            document_id: Document identifier.
            filename: Original filename.

        Returns:
            List of (page_number, png_bytes, ImageInfo).
        """
        import fitz

        doc = fitz.open(str(pdf_path))
        page_numbers = list(range(1, len(doc) + 1))
        doc.close()

        return self.extract_pages(pdf_path, page_numbers, document_id, filename)


# =============================================================================
# S3ImageStore - Upload images to S3
# =============================================================================


class S3ImageStore:
    """Manage image storage in S3."""

    def __init__(
        self,
        bucket: str,
        prefix: str = "images/pptx/",
        region: str | None = None,
    ):
        """Initialize S3 image store.

        Args:
            bucket: S3 bucket name.
            prefix: S3 key prefix for images.
            region: AWS region (uses default if not specified).
        """
        self.bucket = bucket
        self.prefix = prefix.rstrip("/") + "/"
        self.region = region
        self._client = None

    @property
    def client(self):
        """Lazy-load S3 client."""
        if self._client is None:
            try:
                from rag_assist.config.aws_config import get_s3_client
                self._client = get_s3_client()
            except ImportError:
                import boto3
                self._client = boto3.client("s3", region_name=self.region)
        return self._client

    def upload_image(
        self,
        img_bytes: bytes,
        image_info: ImageInfo,
        content_type: str = "image/png",
    ) -> str:
        """Upload image to S3 and update ImageInfo with S3 URI.

        Args:
            img_bytes: Image content.
            image_info: Image metadata (will be updated with S3 info).
            content_type: MIME type.

        Returns:
            S3 URI (s3://bucket/key).
        """
        s3_key = self._generate_s3_key(image_info)

        self.client.put_object(
            Bucket=self.bucket,
            Key=s3_key,
            Body=img_bytes,
            ContentType=content_type,
            Metadata={
                "document_id": image_info.document_id,
                "slide_number": str(image_info.slide_number),
                "image_type": image_info.image_type,
            },
        )

        s3_uri = f"s3://{self.bucket}/{s3_key}"

        # Update ImageInfo
        image_info.s3_uri = s3_uri
        image_info.s3_key = s3_key

        logger.info(
            "Uploaded image to S3",
            s3_uri=s3_uri,
            slide_number=image_info.slide_number,
            size_bytes=image_info.size_bytes,
        )

        return s3_uri

    def upload_images(
        self, images: list[tuple[int, bytes, ImageInfo]]
    ) -> list[str]:
        """Upload multiple images.

        Args:
            images: List of (slide_number, img_bytes, ImageInfo).

        Returns:
            List of S3 URIs.
        """
        uris = []
        for _, img_bytes, info in images:
            uri = self.upload_image(img_bytes, info)
            uris.append(uri)
        return uris

    def get_presigned_url(self, s3_uri: str, expiration: int = 3600) -> str:
        """Generate presigned URL for image access.

        Args:
            s3_uri: S3 URI (s3://bucket/key).
            expiration: URL expiration in seconds.

        Returns:
            Presigned URL.
        """
        if s3_uri.startswith("s3://"):
            parts = s3_uri[5:].split("/", 1)
            bucket = parts[0]
            key = parts[1] if len(parts) > 1 else ""
        else:
            bucket = self.bucket
            key = s3_uri

        url = self.client.generate_presigned_url(
            "get_object",
            Params={"Bucket": bucket, "Key": key},
            ExpiresIn=expiration,
        )
        return url

    def download_image(self, s3_uri: str) -> bytes:
        """Download image from S3."""
        if s3_uri.startswith("s3://"):
            parts = s3_uri[5:].split("/", 1)
            bucket = parts[0]
            key = parts[1] if len(parts) > 1 else ""
        else:
            bucket = self.bucket
            key = s3_uri

        response = self.client.get_object(Bucket=bucket, Key=key)
        return response["Body"].read()

    def delete_document_images(self, document_id: str) -> int:
        """Delete all images for a document."""
        prefix = f"{self.prefix}{document_id}/"
        response = self.client.list_objects_v2(Bucket=self.bucket, Prefix=prefix)

        if "Contents" not in response:
            return 0

        objects = [{"Key": obj["Key"]} for obj in response["Contents"]]
        if objects:
            self.client.delete_objects(Bucket=self.bucket, Delete={"Objects": objects})

        logger.info("Deleted document images", document_id=document_id, count=len(objects))
        return len(objects)

    def _generate_s3_key(self, image_info: ImageInfo) -> str:
        """Generate S3 key from image metadata."""
        return (
            f"{self.prefix}{image_info.document_id}/"
            f"slide_{image_info.slide_number:03d}.png"
        )


# =============================================================================
# Pipeline Function
# =============================================================================


def extract_pptx_images(
    pptx_path: str | Path,
    s3_bucket: str,
    s3_prefix: str = "images/pptx/",
    dpi: int = 150,
    cleanup_pdf: bool = True,
    slide_numbers: list[int] | None = None,
) -> ExtractionResult:
    """Full pipeline: detect → convert → extract → upload.

    Args:
        pptx_path: Path to PPTX file.
        s3_bucket: S3 bucket for image storage.
        s3_prefix: S3 key prefix.
        dpi: Image resolution (higher = better quality).
        cleanup_pdf: Delete intermediate PDF after extraction.
        slide_numbers: Specific slides to process (None = auto-detect).

    Returns:
        ExtractionResult with all extracted images.
    """
    pptx_path = Path(pptx_path)
    filename = pptx_path.name
    errors = []

    logger.info("Starting PPTX image extraction", filename=filename)

    # Step 1: Detect slides with images
    detector = PPTXImageDetector()
    mappings = detector.detect_image_slides(pptx_path)
    document_id = detector._generate_document_id(pptx_path)

    # Get total slides
    prs = Presentation(str(pptx_path))
    total_slides = len(prs.slides)

    # Filter to specific slides if requested
    if slide_numbers:
        mappings = [m for m in mappings if m.slide_number in slide_numbers]

    if not mappings:
        logger.info("No slides with images found", filename=filename)
        return ExtractionResult(
            document_id=document_id,
            filename=filename,
            total_slides=total_slides,
            slides_with_images=0,
        )

    detected_slides = [m.slide_number for m in mappings]
    logger.info("Detected slides with images", slides=detected_slides)

    # Step 2: Convert PPTX to PDF
    converter = PPTXtoPDFConverter()
    pdf_path = None
    try:
        pdf_path = converter.convert(pptx_path)
    except RuntimeError as e:
        errors.append(str(e))
        logger.error("PDF conversion failed", error=str(e))
        return ExtractionResult(
            document_id=document_id,
            filename=filename,
            total_slides=total_slides,
            slides_with_images=len(mappings),
            mappings=mappings,
            errors=errors,
        )

    # Step 3: Extract detected pages as PNG
    extractor = PDFPageExtractor(dpi=dpi)
    try:
        extracted = extractor.extract_pages(
            pdf_path,
            detected_slides,
            document_id,
            filename,
        )
    except Exception as e:
        errors.append(f"Page extraction failed: {str(e)}")
        logger.error("Page extraction failed", error=str(e))
        extracted = []

    # Step 4: Upload to S3
    store = S3ImageStore(s3_bucket, s3_prefix)
    try:
        store.upload_images(extracted)
    except Exception as e:
        errors.append(f"S3 upload failed: {str(e)}")
        logger.error("S3 upload failed", error=str(e))

    # Update mappings with extracted images
    slide_to_info = {page_num: info for page_num, _, info in extracted}
    for mapping in mappings:
        if mapping.slide_number in slide_to_info:
            mapping.images = [slide_to_info[mapping.slide_number]]

    # Cleanup PDF
    if cleanup_pdf and pdf_path and pdf_path.exists():
        try:
            pdf_path.unlink()
            # Also try to remove the temp directory
            if pdf_path.parent.name.startswith("pptx_to_pdf_"):
                pdf_path.parent.rmdir()
        except Exception:
            pass

    # Build result
    all_images = [info for _, _, info in extracted]

    result = ExtractionResult(
        document_id=document_id,
        filename=filename,
        total_slides=total_slides,
        slides_with_images=len(mappings),
        images=all_images,
        mappings=mappings,
        errors=errors,
    )

    logger.info(
        "PPTX image extraction complete",
        filename=filename,
        images_extracted=len(all_images),
        errors=len(errors),
    )

    return result
